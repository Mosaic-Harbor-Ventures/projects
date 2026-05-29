from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

PINK   = RGBColor(0xE8, 0x56, 0x9A)
PINK_L = RGBColor(0xF5, 0xA0, 0xC8)
GREEN  = RGBColor(0x95, 0xC9, 0x3D)
BLUE   = RGBColor(0x4B, 0xC8, 0xE8)
ORANGE = RGBColor(0xF5, 0xA6, 0x23)
DARK   = RGBColor(0x1E, 0x1E, 0x2E)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
MUTED  = RGBColor(0x6B, 0x6B, 0x80)
TEXT   = RGBColor(0x2D, 0x2D, 0x3A)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

SCREENSHOTS = os.path.join(os.path.dirname(__file__), "screenshots")

def blank_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, x, y, w, h, fill_color=None, line_color=None, radius=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        x, y, w, h
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h, size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_pill(slide, text, x, y, w, h, bg_color, text_color=WHITE, size=11):
    r = add_rect(slide, x, y, w, h, fill_color=bg_color)
    add_text(slide, text, x, y, w, h, size=size, bold=True,
             color=text_color, align=PP_ALIGN.CENTER)
    return r

# ─────────────────────────────────────────────
# SLIDE 1 — CAPA
# ─────────────────────────────────────────────
s1 = blank_slide(prs)
bg(s1, DARK)

# decorative circle top-left
c1 = s1.shapes.add_shape(9, Inches(-1), Inches(-1), Inches(5), Inches(5))
c1.fill.solid(); c1.fill.fore_color.rgb = RGBColor(0x3D, 0x10, 0x25)
c1.line.fill.background()

# decorative circle bottom-right
c2 = s1.shapes.add_shape(9, Inches(10), Inches(5), Inches(4), Inches(4))
c2.fill.solid(); c2.fill.fore_color.rgb = RGBColor(0x10, 0x20, 0x3A)
c2.line.fill.background()

# Brand name
add_text(s1, "TRÍADE TEA", Inches(3.5), Inches(0.8), Inches(6.5), Inches(1),
         size=14, bold=True, color=RGBColor(0xAA, 0xAA, 0xBB),
         align=PP_ALIGN.CENTER)

add_text(s1, "CLÍNICA MULTIDISCIPLINAR ABA",
         Inches(3.5), Inches(1.35), Inches(6.5), Inches(0.5),
         size=9, bold=True, color=RGBColor(0x77, 0x77, 0x99),
         align=PP_ALIGN.CENTER)

# divider line
line_box = add_rect(s1, Inches(5.9), Inches(2.0), Inches(1.5), Inches(0.07),
                    fill_color=PINK)

# Title
add_text(s1, "Proposta de Site Profissional",
         Inches(2), Inches(2.3), Inches(9.5), Inches(1),
         size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s1, "para a Tríade TEA",
         Inches(2), Inches(3.2), Inches(9.5), Inches(0.8),
         size=28, bold=True, color=PINK_L, align=PP_ALIGN.CENTER)

# Meta pills row
pill_data = [
    ("🎨  Design + Desenvolvimento", PINK),
    ("📱  Responsivo & SEO", GREEN),
    ("⏱️  Entrega em 7 dias", BLUE),
]
px = Inches(1.5)
for label, col in pill_data:
    add_pill(s1, label, px, Inches(4.4), Inches(3.2), Inches(0.52),
             bg_color=col, size=12)
    px += Inches(3.4)

# Payment note
add_text(s1, "Pagamento em 2x · 50% entrada + 50% na entrega",
         Inches(2), Inches(5.15), Inches(9.5), Inches(0.5),
         size=12, color=RGBColor(0xAA, 0xAA, 0xCC), align=PP_ALIGN.CENTER)

# Proposer
add_text(s1, "Proposta elaborada por Matheus Puppe · 2026",
         Inches(2), Inches(6.7), Inches(9.5), Inches(0.4),
         size=10, color=RGBColor(0x55, 0x55, 0x77), align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# SLIDE 2 — O QUE INCLUI
# ─────────────────────────────────────────────
s2 = blank_slide(prs)
bg(s2, WHITE)

add_pill(s2, "O QUE ESTÁ INCLUSO", Inches(4.9), Inches(0.3), Inches(3.5), Inches(0.4),
         bg_color=GREEN, size=10)
add_text(s2, "Tudo que seu site precisa para converter pacientes",
         Inches(1), Inches(0.85), Inches(11.5), Inches(0.7),
         size=26, bold=True, color=TEXT, align=PP_ALIGN.CENTER)

cards = [
    ("01", "🎨", "Design Personalizado",
     "Layout exclusivo nas cores e identidade da Tríade TEA", PINK),
    ("02", "📱", "Responsivo 100%",
     "Funciona em celular, tablet e computador perfeitamente", GREEN),
    ("03", "📋", "5 Seções Completas",
     "Hero, serviços, sobre, depoimentos e agendamento", BLUE),
    ("04", "💬", "WhatsApp Integrado",
     "Botão flutuante + link direto para conversa na clínica", ORANGE),
    ("05", "🔍", "SEO Básico",
     "Otimização para aparecer no Google e atrair famílias", PINK),
    ("06", "🛠️", "Suporte 30 dias",
     "30 dias de suporte + até 3 revisões incluídas", GREEN),
]

cols = 3
cw = Inches(3.8)
ch = Inches(2.3)
mx = Inches(0.55)
my = Inches(1.8)
gx = Inches(0.3)
gy = Inches(0.25)

for i, (num, icon, title, desc, color) in enumerate(cards):
    col = i % cols
    row = i // cols
    cx = mx + col * (cw + gx)
    cy = my + row * (ch + gy)
    card = add_rect(s2, cx, cy, cw, ch,
                    fill_color=RGBColor(0xF8, 0xF8, 0xFC))
    card.line.color.rgb = RGBColor(0xEE, 0xEE, 0xF5)
    card.line.width = Pt(1)
    # accent bar top
    bar = add_rect(s2, cx, cy, cw, Inches(0.06), fill_color=color)
    add_text(s2, f"{icon}  {title}", cx + Inches(0.2), cy + Inches(0.15),
             cw - Inches(0.4), Inches(0.5), size=14, bold=True, color=TEXT)
    add_text(s2, desc, cx + Inches(0.2), cy + Inches(0.65),
             cw - Inches(0.4), Inches(1.2), size=11, color=MUTED, wrap=True)

# ─────────────────────────────────────────────
# SLIDE 3 — PRÉVIA DO SITE (screenshots)
# ─────────────────────────────────────────────
s3 = blank_slide(prs)
bg(s3, DARK)

add_pill(s3, "PRÉVIA DO SITE", Inches(5.2), Inches(0.25), Inches(3), Inches(0.4),
         bg_color=PINK, size=10)
add_text(s3, "Veja como vai ficar", Inches(1), Inches(0.8), Inches(11.5), Inches(0.7),
         size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

hero_path = os.path.join(SCREENSHOTS, "hero.png")
svc_path  = os.path.join(SCREENSHOTS, "services.png")
age_path  = os.path.join(SCREENSHOTS, "agendamento.png")
dep_path  = os.path.join(SCREENSHOTS, "depoimentos.png")

if os.path.exists(hero_path):
    s3.shapes.add_picture(hero_path, Inches(0.4), Inches(1.7), Inches(7.5), Inches(4.8))
    add_text(s3, "🏠 Página Inicial", Inches(0.4), Inches(6.55), Inches(7.5), Inches(0.4),
             size=11, bold=True, color=RGBColor(0xAA, 0xAA, 0xBB))

if os.path.exists(svc_path):
    s3.shapes.add_picture(svc_path, Inches(8.1), Inches(1.7), Inches(5.0), Inches(2.3))
    add_text(s3, "🩺 Serviços", Inches(8.1), Inches(4.05), Inches(5.0), Inches(0.35),
             size=10, bold=True, color=RGBColor(0xAA, 0xAA, 0xBB))

if os.path.exists(age_path):
    s3.shapes.add_picture(age_path, Inches(8.1), Inches(4.5), Inches(5.0), Inches(2.0))
    add_text(s3, "📅 Agendamento", Inches(8.1), Inches(6.55), Inches(5.0), Inches(0.35),
             size=10, bold=True, color=RGBColor(0xAA, 0xAA, 0xBB))

# ─────────────────────────────────────────────
# SLIDE 4 — COMO FUNCIONA
# ─────────────────────────────────────────────
s4 = blank_slide(prs)
bg(s4, RGBColor(0xF8, 0xF4, 0xFF))

add_pill(s4, "COMO FUNCIONA", Inches(5.2), Inches(0.3), Inches(3), Inches(0.4),
         bg_color=BLUE, size=10)
add_text(s4, "Do acordo à entrega em 7 dias",
         Inches(1), Inches(0.85), Inches(11.5), Inches(0.7),
         size=26, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
add_text(s4, "Processo simples, transparente e rápido. Você foca na clínica, eu cuido do site.",
         Inches(2), Inches(1.55), Inches(9.5), Inches(0.5),
         size=13, color=MUTED, align=PP_ALIGN.CENTER)

steps = [
    ("🤝", "1", "Aprovação", "Proposta aceita e pagamento da entrada de R$ 1.250", PINK),
    ("📋", "2", "Briefing", "Coleto informações, textos, fotos e detalhes da clínica", ORANGE),
    ("💻", "3", "Desenvolvimento", "Site desenvolvido em até 5 dias com atualizações", GREEN),
    ("🚀", "4", "Entrega", "Aprovação final, pagamento e site no ar", BLUE),
]

sw = Inches(2.8)
sx_start = Inches(0.6)
sy = Inches(2.4)
for i, (icon, num, title, desc, col) in enumerate(steps):
    sx = sx_start + i * (sw + Inches(0.25))
    circ = s4.shapes.add_shape(9, sx + Inches(0.8), sy, Inches(1.2), Inches(1.2))
    circ.fill.solid(); circ.fill.fore_color.rgb = col; circ.line.fill.background()
    add_text(s4, icon, sx + Inches(0.8), sy, Inches(1.2), Inches(1.2),
             size=26, align=PP_ALIGN.CENTER, color=WHITE)
    add_text(s4, title, sx, sy + Inches(1.35), sw, Inches(0.45),
             size=14, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
    add_text(s4, desc, sx, sy + Inches(1.85), sw, Inches(1.0),
             size=11, color=MUTED, align=PP_ALIGN.CENTER, wrap=True)

# badges row
badge_data = [
    ("⏱️ 7 dias", "Prazo entrega", PINK),
    ("🔄 3 revisões", "Ajustes incluídos", GREEN),
    ("🛡️ 30 dias", "Suporte pós-entrega", BLUE),
    ("💳 2x", "Pagamento parcelado", ORANGE),
]
bw = Inches(2.8)
bx_start = Inches(0.6)
for i, (val, lbl, col) in enumerate(badge_data):
    bx = bx_start + i * (bw + Inches(0.25))
    card = add_rect(s4, bx, Inches(5.6), bw, Inches(1.3),
                    fill_color=WHITE)
    card.line.color.rgb = RGBColor(0xDD, 0xDD, 0xEE); card.line.width = Pt(1)
    add_text(s4, val, bx, Inches(5.7), bw, Inches(0.55),
             size=16, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(s4, lbl, bx, Inches(6.25), bw, Inches(0.4),
             size=11, color=MUTED, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# SLIDE 5 — INVESTIMENTO
# ─────────────────────────────────────────────
s5 = blank_slide(prs)
bg(s5, WHITE)

add_pill(s5, "INVESTIMENTO", Inches(5.2), Inches(0.3), Inches(3), Inches(0.4),
         bg_color=ORANGE, size=10)
add_text(s5, "Um site que trabalha por você 24h por dia",
         Inches(1), Inches(0.85), Inches(11.5), Inches(0.7),
         size=26, bold=True, color=TEXT, align=PP_ALIGN.CENTER)

# Dark pricing box
box = add_rect(s5, Inches(1.0), Inches(1.8), Inches(11.5), Inches(4.8),
               fill_color=DARK)

# left column — includes
left_items = [
    "✓  Design personalizado Tríade TEA",
    "✓  5 seções completas",
    "✓  Responsivo (mobile + tablet + desktop)",
    "✓  WhatsApp integrado",
    "✓  Formulário de agendamento",
    "✓  SEO básico configurado",
    "✓  3 revisões incluídas",
    "✓  30 dias de suporte",
]
add_text(s5, "Pacote Completo — Tríade TEA",
         Inches(1.3), Inches(2.05), Inches(5.5), Inches(0.55),
         size=16, bold=True, color=WHITE)
for j, item in enumerate(left_items):
    add_text(s5, item, Inches(1.3), Inches(2.75) + Inches(j * 0.38),
             Inches(5.5), Inches(0.4), size=12, color=RGBColor(0xCC, 0xCC, 0xDD))

# right column — price
add_text(s5, "R$", Inches(7.8), Inches(2.4), Inches(1), Inches(0.8),
         size=22, bold=True, color=PINK_L, align=PP_ALIGN.RIGHT)
add_text(s5, "2.500", Inches(8.6), Inches(2.0), Inches(3.5), Inches(1.5),
         size=72, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# payment box
pay_box = add_rect(s5, Inches(7.8), Inches(3.65), Inches(4.2), Inches(1.5),
                   fill_color=RGBColor(0x2E, 0x2E, 0x42))
pay_box.line.color.rgb = RGBColor(0x44, 0x44, 0x60); pay_box.line.width = Pt(1)
add_text(s5, "FORMA DE PAGAMENTO",
         Inches(7.85), Inches(3.75), Inches(4.1), Inches(0.35),
         size=9, bold=True, color=RGBColor(0x77, 0x77, 0x99), align=PP_ALIGN.CENTER)
add_text(s5, "🟡  Entrada (aprovação)       R$ 1.250",
         Inches(7.85), Inches(4.15), Inches(4.1), Inches(0.35),
         size=12, color=ORANGE, align=PP_ALIGN.CENTER)
add_text(s5, "🟢  Final (entrega)              R$ 1.250",
         Inches(7.85), Inches(4.55), Inches(4.1), Inches(0.35),
         size=12, color=GREEN, align=PP_ALIGN.CENTER)

cta = add_rect(s5, Inches(7.8), Inches(5.35), Inches(4.2), Inches(0.65),
               fill_color=PINK)
cta.line.fill.background()
add_text(s5, "💬  Quero aprovar a proposta",
         Inches(7.8), Inches(5.35), Inches(4.2), Inches(0.65),
         size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s5, "Pix · Transferência · Cartão de crédito",
         Inches(7.8), Inches(6.1), Inches(4.2), Inches(0.35),
         size=10, color=RGBColor(0x88, 0x88, 0xAA), align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# SLIDE 6 — ENCERRAMENTO
# ─────────────────────────────────────────────
s6 = blank_slide(prs)
bg(s6, DARK)

# glow circle
glow = s6.shapes.add_shape(9, Inches(4), Inches(1), Inches(5.5), Inches(5.5))
glow.fill.solid(); glow.fill.fore_color.rgb = RGBColor(0x35, 0x10, 0x25)
glow.line.fill.background()

add_text(s6, "🚀", Inches(5.9), Inches(0.6), Inches(1.5), Inches(1),
         size=40, align=PP_ALIGN.CENTER, color=WHITE)

add_text(s6, "Pronto para colocar",
         Inches(1), Inches(1.6), Inches(11.5), Inches(0.7),
         size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s6, "a Tríade TEA no mapa digital?",
         Inches(1), Inches(2.25), Inches(11.5), Inches(0.7),
         size=30, bold=True, color=PINK_L, align=PP_ALIGN.CENTER)

add_text(s6,
         "Famílias buscam clínicas de TEA online todos os dias.\nUm site profissional coloca a Tríade TEA na frente da concorrência.",
         Inches(2), Inches(3.1), Inches(9.5), Inches(0.9),
         size=14, color=RGBColor(0xAA, 0xAA, 0xCC), align=PP_ALIGN.CENTER, wrap=True)

# chips
chips = [
    ("R$ 2.500 · 2x", PINK),
    ("⏱️ Entrega em 7 dias", BLUE),
    ("matheus.puppe@gmail.com", GREEN),
]
cx_chip = Inches(1.5)
for label, col in chips:
    chip = add_rect(s6, cx_chip, Inches(4.3), Inches(3.3), Inches(0.55),
                    fill_color=RGBColor(0x2E, 0x2E, 0x44))
    chip.line.color.rgb = col; chip.line.width = Pt(1.5)
    add_text(s6, label, cx_chip, Inches(4.3), Inches(3.3), Inches(0.55),
             size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    cx_chip += Inches(3.5)

add_text(s6, "Proposta elaborada por Matheus Puppe · 2026",
         Inches(1), Inches(6.8), Inches(11.5), Inches(0.4),
         size=10, color=RGBColor(0x44, 0x44, 0x66), align=PP_ALIGN.CENTER)

# SAVE
out = os.path.join(os.path.dirname(__file__), "triade_tea_proposta.pptx")
prs.save(out)
print(f"Saved: {out}")
